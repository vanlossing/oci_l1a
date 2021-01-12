#!/usr/bin/env python
"""test read_oci_1a_class"""

__version__ = "0.0.1"
__author__ = "Christopher Field <Christopher.T.Field@nasa.gov>"

import argparse
from copy import deepcopy
import datetime
from math import ceil, floor
from pprint import pprint
import os
from pathlib import Path
import sys

from matplotlib import pyplot as plt
import matplotlib as mpl                # Used to get color map.
from mpl_toolkits.mplot3d import axes3d
from netCDF4 import Dataset         #https://unidata.github.io/netcdf4-python/netCDF4/index.html
import numpy as np
import numpy.ma as ma
from pptx import Presentation
from pptx.util import Inches
from scipy.fftpack import fft
from scipy.optimize import curve_fit

# sys.path.append(os.path.abspath(os.curdir))
# import .context

from python_pptx_helper.create_ppt import create_ppt
import read_oci_1a_class as oci_ss
# from read_oci_1a_class import READ_OCI_SPA_SPE_L1A as oci_ss

def commandline():
    """Define the command line arguments. Return parser.

    Usage: args = commandline().parse_args()

    This structure permits script documentation by sphnix. Use  .. argparse::.
    See https://readthedocs.org/projects/sphinx-argparse/ for more info.
    """

    parser = argparse.ArgumentParser(
        description=__doc__,
        epilog=" "
    )
    parser.add_argument('input',
                        help='Input argument')

    return parser


def main(args):
    """Main routine. Could be imported and called by another Python program

    :param args:    Argparse generated command line arguments
    :return:        None
    """
    pass


def convert_sod2datetime(x, origin):
    """Convert the list of x (seconds of day) to array of seconds of day and datetimes

    :param x: list of seconds of day (float)
    :param origin: Datetime object of midnight that starts the day.
    :return: time_s, time_dt. 1darray of seconds since midnight and datetimes"""
    time_s = np.asarray(x)
    time_dt = np.asarray([origin + datetime.timedelta(seconds=s) for s in x])
    return time_s, time_dt

def func(x, a, b, c, d):
    """Function a * (1-exp(-(x-b)/c) + d to fit

    :param x: array_like times (floats) at which to compute the function value
    :param a: scalar step size
    :param b: scalar time of step start
    :param c: scalar response time constant (same units as x)
    :param d: scalar starting value for step response.
    :return: function evaluated at each x value
    """
    return a * (1 - np.exp(-(x - b)/c)) + d


def plot_time_history(dark, channels, ylabels, filename, start_time, figsize, gridspec_kw):

    print("In plot_time_history: ", filename)
    gridspec_kw = {"left" : 0.07,
                   "right" : 0.95,
                   "top" : 0.90,
                   "bottom" : 0.15,
                   "wspace" : 0.10,
                   "hspace" : 0.20,
                   # "width_ratios" : [1],
                   # "height_ratios" : [1, 1, 1],
                   "width_ratios" : [4, 1],
                   "height_ratios" : [1, 1, 1],
                   }

    fig_5s, ax_5s = plt.subplots( len(channels), 2, sharex="col", sharey="row", squeeze=False,
                                  figsize=figsize, gridspec_kw=gridspec_kw)
    fig_1m, ax_1m = plt.subplots( len(channels), 2, sharex="col", sharey="row",  squeeze=False,
                                  figsize=figsize, gridspec_kw=gridspec_kw)
    fig_5s.suptitle("5 Second Sample starting at {:}".format(start_time))
    fig_1m.suptitle('"1 Minute" Sample starting at {:}'.format(start_time))

    fig_fft, ax_fft = plt.subplots(len(channels), 1, sharex=True, figsize=figsize) #, gridspec_kw=gridspec_kw)
    fig_fft.suptitle("FFT of 1 Minute Sample starting at {:}".format(start_time))

    for chan_idx, (chan, ylabel)  in enumerate(zip(channels, ylabels)):
        data = dark[:, chan, :]
        data2 = data.ravel(order='C')               # Keep the spatial pixels of each scan together
        ax_5s[chan_idx, 0].set_title(ylabel)
        ax_5s[chan_idx, 0].plot(data2[:1500], '-', linewidth=0.25)
        ax_5s[chan_idx, 0].set_xlim(xmin=0)
        ylim1 = ax_5s[chan_idx, 0].get_ylim()
        ylim1 = (floor(ylim1[0])-0.5, ceil(ylim1[1])+0.5)
        bins1 = int(ylim1[1] - ylim1[0])
        data3 = data2.compressed()
        ax_5s[chan_idx, 1].hist(data3[:1500],
                                bins=bins1,
                                range=ylim1,
                                orientation="horizontal")
        ax_5s[chan_idx, 1].text(0.80, 0.9, '$\sigma$ = {:.1f}'.format(data2[:1500].std()),
                                horizontalalignment='center', verticalalignment = 'center', transform = ax_5s[chan_idx, 1].transAxes)
        ax_1m[chan_idx, 0].set_title(ylabel)
        ax_1m[chan_idx, 0].plot(data2, '-', linewidth=0.25)
        ax_1m[chan_idx, 0].set_xlim(xmin=0)
        ylim2 = ax_1m[chan_idx, 0].get_ylim()
        ylim2 = (floor(ylim2[0]), ceil(ylim2[1]))
        bins2 = int(ylim2[1] - ylim2[0])
        ax_1m[chan_idx, 1].hist(data3,
                                bins=bins2,
                                range=ylim2,
                                orientation="horizontal")
        ax_1m[chan_idx, 1].text(0.80, 0.9, '$\sigma$ = {:.1f}'.format(data2.std()),
                                horizontalalignment='center', verticalalignment = 'center', transform = ax_1m[chan_idx, 1].transAxes)

        data3 = data2.filled(fill_value=data2.mean())
        time_series = data3 - data3.mean()
        power = np.abs(fft(time_series))
        ax_fft[chan_idx].plot(power[:power.shape[0]//2], '-', linewidth=0.25)
        y_min, y_max = ax_fft[chan_idx].get_ylim()
        ax_fft[chan_idx].set_ylim(0, 1000)
        # if y_max > 30000:
        #     ax_fft[chan_idx].set_ylim(0, 30000)
        ax_fft[chan_idx].set_title(ylabel)

    fig_5s.savefig(filename+"_5s.png", dpi=1200)
    fig_1m.savefig(filename+"_1m.png", dpi=1200)
    # fig_fft.savefig(filename+"_fft.png", dpi=1200)

    del fig_5s, fig_1m
    fig_5s =  ax_5s = fig_1m = ax_1m = None
    return ((fig_5s, ax_5s), (fig_1m, ax_1m)), filename
    # return ((fig_5s, ax_5s), (fig_1m, ax_1m), (fig_fft, ax_fft)), filename


def plot_response(x, y_mean, y_std, ax, ylabel, origin, plot_stable=True):

    xs, xdt = convert_sod2datetime(x, origin)
    print("First data point is at {:}    {:}".format(xs[0], xdt[0]))
    ax[0].set_ylabel(ylabel)
    ax[0].plot(xdt, y_mean, 'x')
    ax[1].plot(xdt, y_std, 'x')
    try:
        popt, pcov = curve_fit(func, xs, y_mean, p0=(y_mean.max() - y_mean.min(),
                                                     x[0],
                                                     3600.,
                                                     y_mean.min()))
    except RuntimeError:
        return                                  # Didn't get a fit, done.
    xs_p, xdt_p = convert_sod2datetime(np.linspace(x[0], x[-1], 20), origin)
    print("Fit solution is ", popt)
    ax[0].plot(xdt_p, func(xs_p, *popt), 'r-')
    ax[0].text(0.20, 0.1,
               "{0:.0f} (1-exp[ -t / {2:.0f} s ]) + {3:.0f}   Time origin is {4:}"
                    .format(*popt,  convert_sod2datetime([popt[1]], origin)[1][0].strftime("%H:%M")),
               transform=ax[0].transAxes)
    y = func(np.array([4800, 5400])+x[0], *popt)               # Compute y at 50 and 60 minutes

    if not plot_stable:
        return                                  # Don't need the stable limits
    t_10 = -popt[2] * np.log(0.100 * (popt[0]+popt[3]) / popt[0] / (1.100 - np.exp(-600/popt[2]))) + popt[1]
    t_1  = -popt[2] * np.log(0.010 * (popt[0]+popt[3]) / popt[0] / (1.010 - np.exp(-600/popt[2]))) + popt[1]
    t_05 = -popt[2] * np.log(0.005 * (popt[0]+popt[3]) / popt[0] / (1.005 - np.exp(-600/popt[2]))) + popt[1]
    t_01 = -popt[2] * np.log(0.001 * (popt[0]+popt[3]) / popt[0] / (1.001 - np.exp(-600/popt[2]))) + popt[1]

    print("To 1%.  {:.0f} seconds = {:}"
          .format(t_1, (origin + datetime.timedelta(seconds=t_1)).strftime("%y-%m-%d %H:%M")))
    print("To 0.5%  {:.0f} seconds = {:}"
          .format(t_05, (origin + datetime.timedelta(seconds=t_05)).strftime("%y-%m-%d %H:%M")))
    print("To 0.1%  {:.0f} seconds = {:}"
          .format(t_01, (origin + datetime.timedelta(seconds=t_01)).strftime("%y-%m-%d %H:%M")))
    print("Stability over 10 minutes after 1 hour is {:.1f}%".format((y[1]-y[0])/y[0] * 100))
    print(origin + datetime.timedelta(seconds=t_10),
          origin + datetime.timedelta(seconds=t_1),
          origin + datetime.timedelta(seconds=t_01) )
    xs, xdt = convert_sod2datetime([popt[1], t_1-600, t_1, t_05-600, t_05, t_01-600, t_01], origin)
    y_xs = func(xs, *popt)
    # [ax[0].axvline(x) for x in xdt]
    # [ax[0].axhline(y) for y in y_xs]
    ax[0].plot([xdt[1], xdt[2], xdt[2]], [y_xs[1], y_xs[1], y_xs[2]], 'k')
    ax[0].plot([xdt[3], xdt[4], xdt[4]], [y_xs[3], y_xs[3], y_xs[4]], 'k')
    ax[0].plot([xdt[5], xdt[6], xdt[6]], [y_xs[5], y_xs[5], y_xs[6]], 'k')
    ax[0].text(xdt[2], y_xs[1], " 1%")
    ax[0].text(xdt[4], y_xs[3], " 0.5%")
    ax[0].text(xdt[6], y_xs[5], " 0.1%")
    ax[0].text(0.20, 0.2,
               "10 minutes stability 1.5 hour after data start is {:.1f}%".format((y[1]-y[0])/y[1] * 100),
               transform=ax[0].transAxes)


# def create_ppt(input, output, populate, show_placeholders=False):
#     """Use the input powerpoint file as the template for the output powerpoint file.
#
#     :param input: Path to source pptx file for templates. If None, use internal default template file.
#     :param output: Path to output pptx file. If the same as source, will overwrite input without warning.
#     :param populate: A list of values to populate the slides. See below.
#     :param show_placeholders: If True, don't create any slides but print structure of selected slides.
#     :return: None.
#
#     Taken from https://pbpython.com/creating-powerpoint.html on 3 November 2020
#
#     populate is a list of dictionaries of data to populate the slides. Each element has the following format::
#
#         - "layout" : number of layout to select
#         - "title" : String for the title, if needed.
#         - "text" : List of (placeholder #, text) for the slide.
#         - "picture" : List of (placeholder #, path to picture file) for the slide.
#
#     If show_placeholders is True, then only the "layout" is required for each slide. The rest of the
#     dictionary entries are ignored.
#     """
#     prs = Presentation(input)
#
#     for slide_num, pop_slide in enumerate(populate):
#         slide_layout = prs.slide_layouts[pop_slide["layout"]]
#         if show_placeholders:
#             print("Slide {:}".format(slide_num))
#             for shape in slide_layout.placeholders:
#                 print(shape.placeholder_format.idx, shape.name, shape.shape_type)
#             continue                            # Don't populate the slide
#
#         slide = prs.slides.add_slide(slide_layout)
#         try:                                    # Process all of the text boxes
#             for idx, text in pop_slide["text"]:
#                 slide.placeholders[idx].text = text
#         except KeyError:
#             pass                                # Didn't request filling any text placeholders.
#
#         try:                                    # Process all of the picture boxes
#             for idx, path in pop_slide["picture"]:
#                 slide.placeholders[idx].insert_picture(path)
#         except KeyError:
#             pass                                # Didn't request filling any text placeholders.
#     if (not show_placeholders) and (output is not None):
#         prs.save(output)
#     return


if __name__ == "__main__":
    ppx_input =  "/Users/cfield/Documents/PACE/python_code/Presentation_template_in2.pptx"
    ppx_output = "/Users/cfield/Documents/PACE/python_code/OCI_REPORT_2020__startup_transients.pptx"
    if 0:
        data_path = Path("/Users/cfield/Documents/PACE/Data/Cold-Start-Test/2020-10-29")
        midnight = datetime.datetime(2020, 10, 29, 0, 0, 0)
    else:
        data_path = Path("/Users/cfield/Documents/PACE/Data/Cold-Start-Test/2020-10-30")
        midnight = datetime.datetime(2020, 10, 30, 0, 0, 0)
        ppx_input = ppx_output                  # Append new slides to old file

    # ppx_output = None
    pattern = ""
    dark_zero_limit = 40
    zero_limit = 41

    figsize = (11.5, 5.75)
    gridspec_kw = {"left" : 0.07,
                   "right" : 0.95,
                   "top" : 0.90,
                   "bottom" : 0.15,
                   "wspace" : 0.10,
                   "hspace" : 0.20
                   }

    files = oci_ss.get_pace_oci_l1a_files(data_path, pattern)
    pprint(files)

    # oci = oci_ss.read_oci_1a_class(None)
    # for file in path:
    #     oci.append(oci_ss.read_oci_1a_class(file))
    # oci.append_conclude()

    time_lst = []
    Sci_red_610_mean_lst, Sci_red_610_std_lst = [], []
    Sci_red_685_mean_lst, Sci_red_685_std_lst = [], []
    Sci_red_885_mean_lst, Sci_red_885_std_lst = [], []
    Sci_swir_mean_lst, Sci_swir_std_lst = [], []
    # Thermal lists
    time_thermal_lst = []
    aob_temp_lst, dau_temp_lst = [], []
    red_fpa_temp_lst, sds_det_temp_lst = [], []
    start_time = 0                                  # Used to generate 1 minute plots every hour
    time_series_names_lst = []

    for file_idx, file in enumerate(files):
        # Collect the mean and std CCD and SWIR data into lists
        oci = oci_ss.READ_OCI_SPA_SPE_L1A(file)

        i_610 = np.flatnonzero(oci.red_wavelengths == 610)
        i_685 = np.flatnonzero(oci.red_wavelengths == 685)
        i_885 = np.flatnonzero(oci.red_wavelengths == 885)

        time_lst.append(oci.scan_start_time[0])
        Sci_red_610_mean = oci.DC_red_data[1:-1, i_610, :].mean()
        Sci_red_610_std = oci.DC_red_data[1:-1, i_610, :].std()
        Sci_red_610_mean_lst.append(Sci_red_610_mean), Sci_red_610_std_lst.append(Sci_red_610_std)

        Sci_red_685_mean = oci.DC_red_data[1:-1, i_685, :].mean()
        Sci_red_685_std = oci.DC_red_data[1:-1, i_685, :].std()
        Sci_red_685_mean_lst.append(Sci_red_685_mean), Sci_red_685_std_lst.append(Sci_red_685_std)

        Sci_red_885_mean = oci.DC_red_data[1:-1, i_885, :].mean()
        Sci_red_885_std = oci.DC_red_data[1:-1, i_885, :].std()
        Sci_red_885_mean_lst.append(Sci_red_885_mean), Sci_red_885_std_lst.append(Sci_red_885_std)

        mean_lst, std_lst = [], []              # Temp lists to gather SWIR channels of interest
        for i in [0, 1, 3]:
            mean_lst.append(oci.DC_swir_data[1:-1, i, :].mean())
            std_lst.append(oci.DC_swir_data[1:-1, i, :].std())
        Sci_swir_mean_lst.append(mean_lst), Sci_swir_std_lst.append(std_lst)

        if oci.scan_start_time[1] - start_time > 3600:
            # print("np.diff(oci.scan_start_time)")
            # print(np.diff(oci.scan_start_time))
            # print("1.0 / np.diff(oci.scan_start_time)")
            # print(1.0 / np.diff(oci.scan_start_time))
            # print("oci.scan_start_time[-1] - oci.scan_start_time[0]")
            print(oci.scan_start_time[-2] - oci.scan_start_time[1])
            start_time = oci.scan_start_time[1]
            start_datetime = convert_sod2datetime([start_time], midnight)[1][0]
            fig_R, filenameA = plot_time_history(oci.DC_red_data[1:-1, :, :],
                                                 [i_610, i_685, i_885],
                                                 ["610 nm", "685 nm", "885 nm"],
                                                 "/Users/cfield/Desktop/temp/red_dark_{:d}"
                                                         .format(int(start_time)),
                                                 start_datetime,
                                                 figsize,
                                                 gridspec_kw)
            fig_I, filenameB = plot_time_history(oci.DC_swir_data[1:-1, :, :],
                                                 [0, 1, 3],
                                                 ["940 nm", "1038 nm", "1250 nm (HG)"],
                                                 "/Users/cfield/Desktop/temp/swir_dark_{:d}"
                                                          .format(int(start_time)),
                                                 start_datetime,
                                                 figsize,
                                                 gridspec_kw)
            time_series_names_lst.append((filenameA, filenameB))
            # plt.show()

        # Collect the mean and std temperature data into lists
        oci = oci_ss.READ_OCI_THERMAL_L1A(file)
        time_thermal_lst.append(oci.dau_tlm_time[0])
        aob_temp_lst.append(oci.aob_temperatures.mean(axis=0))
        dau_temp_lst.append(oci.dau_temperatures.mean(axis=0))
        red_fpa_temp_lst.append(oci.red_fpa_temperatures.mean(axis=0))
        sds_det_temp_lst.append(oci.sds_det_temperatures.mean(axis=0))

    time_xs, time_xdt = convert_sod2datetime(time_lst, origin=midnight)

    Sci_red_610_mean = np.asarray(Sci_red_610_mean_lst)
    Sci_red_610_std = np.asarray(Sci_red_610_std_lst)
    Sci_red_685_mean = np.asarray(Sci_red_685_mean_lst)
    Sci_red_685_std = np.asarray(Sci_red_685_std_lst)
    Sci_red_885_mean = np.asarray(Sci_red_885_mean_lst)
    Sci_red_885_std = np.asarray(Sci_red_885_std_lst)

    Sci_swir_mean = np.asarray(Sci_swir_mean_lst)
    Sci_swir_std = np.asarray(Sci_swir_std_lst)

    # Convert temperatures to ndarrays
    time_xs_thermal, time_xdt_thermal = convert_sod2datetime(time_thermal_lst, origin=midnight)
    aob_temp = np.asarray(aob_temp_lst)
    dau_temp = np.asarray(dau_temp_lst)
    red_fpa_temp = np.asarray(red_fpa_temp_lst)
    sds_det_temp = np.asarray(sds_det_temp_lst)

    figsize = (11.5, 5.75)
    gridspec_kw = {"left" : 0.07,
                   "right" : 0.95,
                   "top" : 0.90,
                   "bottom" : 0.15,
                   "wspace" : 0.10,
                   "hspace" : 0.20
                   }

    fig, ax = plt.subplots( 3, 2, sharex=True, figsize=figsize, gridspec_kw=gridspec_kw)
    fig.suptitle("Red Dark Mean and Standard Deviation (DN16) Collected {:}".format(data_path.name))

    ax[0, 0].set_title("Mean Dark Level")
    ax[0, 1].set_title("Standard Deviation Dark Level")
    fig.autofmt_xdate(rotation=45)

    plot_response(time_lst, Sci_red_610_mean, Sci_red_610_std, ax[0], "610 nm", midnight)
    plot_response(time_lst, Sci_red_685_mean, Sci_red_685_std, ax[1], "685 nm", midnight)
    plot_response(time_lst, Sci_red_885_mean, Sci_red_885_std, ax[2], "885 nm", midnight)

    fig_filename1 = "/Users/cfield/Desktop/temp/ccd_transients.png"
    fig.savefig(fig_filename1, dpi=600)

    # Plot the desired SWIR channels
    fig, ax = plt.subplots( 3, 2, sharex=True, figsize=figsize, gridspec_kw=gridspec_kw)
    fig.suptitle("SWIR Dark Mean and Standard Deviation Collected {:}".format(data_path.name))
    ax[0, 0].set_title("Mean Dark Level")
    ax[0, 1].set_title("Standard Deviation Dark Level")
    fig.autofmt_xdate(rotation=45)
    for a, m, s, ylabel in zip(ax,
                               Sci_swir_mean.T,
                               Sci_swir_std.T,
                               ["SWIR 1 (940 nm)","SWIR 2 (1038 nm", "SWIR 3 (1250 nm HG"]):
        plot_response(time_lst, m, s, a, ylabel, midnight, plot_stable=False)

    fig_filename2 = "/Users/cfield/Desktop/temp/swir_transients.png"
    fig.savefig(fig_filename2, dpi=600)

    # Plot the thermal data.
    fig, ax = plt.subplots(4, 2, sharex=True, figsize=figsize, gridspec_kw=gridspec_kw)
    fig.patch.set_facecolor('#E0E0E0')
    fig.suptitle("Mean Thermal Data Collected {:}".format(data_path.name))
    fig.autofmt_xdate(rotation=45)
    l = aob_temp.shape[1] // 2 + 1
    ax[0, 0].set_title("aob_temp Channels 0 - {:}".format(l-1))
    ax[0, 1].set_title("aob_temp Channels {:} to {:}".format(l, aob_temp.shape[1]-1))
    ax[0, 0].plot(time_xdt_thermal, aob_temp[:, :l], "-+")
    ax[0, 1].plot(time_xdt_thermal, aob_temp[:, l:], "-+")
    l = dau_temp.shape[1] // 2
    ax[1, 0].set_title("dau_temp Channels 0 - {:}".format(l-1))
    ax[1, 1].set_title("dau_temp Channels {:} to {:}".format(l, dau_temp.shape[1]-1))
    ax[1, 0].plot(time_xdt_thermal, dau_temp[:, :l], "-+")
    ax[1, 1].plot(time_xdt_thermal, dau_temp[:, l:], "-+")
    l = red_fpa_temp.shape[1] // 2
    ax[2, 0].set_title("red_fpa_temp Channels 0 - {:}".format(l-1))
    ax[2, 1].set_title("red_fpa_temp {:} to {:}".format(l, red_fpa_temp.shape[1]-1))
    ax[2, 0].plot(time_xdt_thermal, red_fpa_temp[:, :l], "-+")
    ax[2, 1].plot(time_xdt_thermal, red_fpa_temp[:, l:], "-+")
    l = sds_det_temp.shape[1] // 2
    ax[3, 0].set_title("sds_det_temp Channels 0 - {:}".format(l-1))
    ax[3, 1].set_title("sds_det_temp {:} to {:}".format(l, sds_det_temp.shape[1]-1))
    ax[3, 0].plot(time_xdt_thermal, sds_det_temp[:, :l], "-+")
    ax[3, 1].plot(time_xdt_thermal, sds_det_temp[:, l:], "-+")

    fig_filename3 = "/Users/cfield/Desktop/temp/thermal_transients.png"
    fig.savefig(fig_filename3, dpi=600)

    print("Processing is all done. Now make the Power Point file.")
    print("Defining the slides")

    slides = [{"layout" : 0,    # Title slide
               "text" : [(0, "OCI Start Up Dark Transients to Thermal Equilibrium"),
                         (1, "Generated on {:%m-%d-%Y}\n by Christopher Field".format(datetime.date.today()))]
               },
              {"layout" : 11,   # Picture slide
               "text" : [(0, "Red Dark Level (DN16)"),
                         # (10, "Date"),
                         # (11, "Footer"),
                         # (12, "Slide #")
                         (13, "Source data from {:}".format(data_path)),
                        ],
               "picture" : [(1, fig_filename1)]
               },
              {"layout" : 11,   # Picture slide
               "text" : [(0, "SWIR Dark Level (DN)"),
                         (13, "Source data from {:}".format(data_path)),
                        ],
               "picture" : [(1, fig_filename2)]
               },
              {"layout" : 11,   # Picture slide
               "text" : [(0, "All Temperature Sensors. I believe the vertical scale is deg C, but not sure."),
                         (13, "Source data from {:}".format(data_path)),
                        ],
               "picture" : [(1, fig_filename3)]
               },
              ]

    for names in time_series_names_lst:                 # Collect the Red Time Series
        slides.append({"layout" : 11,   # Picture slide
                       "text" : [(0, "5 Second Time Series of Red Dark Counts"),
                                 (13, "Source data from {:}".format(data_path)),
                        ],
                       "picture" : [(1, names[0]+"_5s.png")]
                       },
                     )
        slides.append({"layout" : 11,   # Picture slide
                       "text" : [(0, '"1 Minute" Time Series of Red Dark Counts'),
                                 (13, "Source data from {:}".format(data_path)),
                        ],
                       "picture" : [(1, names[0]+"_1m.png")]
                       },
                     )
    for names in time_series_names_lst:                 # Collect the IR Time Series
        slides.append({"layout" : 11,   # Picture slide
                       "text" : [(0, "5 Second Time Series of SWIR Dark Counts"),
                                 (13, "Source data from {:}".format(data_path)),
                        ],
                       "picture" : [(1, names[1]+"_5s.png")]
                       },
                     )
    for names in time_series_names_lst:                 # Collect the IR Time Series
        slides.append({"layout" : 11,   # Picture slide
                       "text" : [(0, "1 Minute Time Series of SWIR Dark Counts"),
                                 (13, "Source data from {:}".format(data_path)),
                        ],
                       "picture" : [(1, names[1]+"_1m.png")]
                       },
                     )
    print("Creating the file.")
    create_ppt(ppx_input, ppx_output, slides, show_placeholders=False)

    plt.show()
